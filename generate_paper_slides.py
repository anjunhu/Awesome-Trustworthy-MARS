"""Per-paper PPTX — no accent bar, +10pt fonts, full data from paper_data.py"""
import re, sys
sys.path.insert(0, "/Users/anjunhu/Bookchapter")
from paper_data import KNOWN
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def parse_readme(path="README.md"):
    text = open(path).read()
    sections = re.split(r'\n## ', text)
    rf_lookup = {
        'Foundational':'RF0','Risk Family 1':'RF1','Risk Family 2':'RF2',
        'Risk Family 3':'RF3','Risk Family 4':'RF4','Risk Family 5':'RF5',
        'Risk Family 6':'RF6','Collusion':'RF6','Fairness':'RF5',
        'Evaluation':'EVAL','Defence':'DEF','Broad Safety':'SURVEY','Uncategorised':'TBD',
    }
    papers = []
    for sec in sections:
        lines = sec.strip().split('\n')
        sec_title = lines[0].strip()
        rf = next((v for k,v in rf_lookup.items() if k in sec_title), 'OTHER')
        for title,authors,venue,arxiv,notes,tags in re.findall(
            r'\|\s*\*\*(.*?)\*\*\s*[—–-]\s*(.*?)\s*\|\s*(.*?)\s*\|\s*\[?([\d]{4}\.\d+)\]?.*?\|\s*(.*?)\s*\|\s*(.*?)\s*\|', sec):
            rfs = re.findall(r'`risk:(rf\d)`', tags)
            topics = ", ".join(re.findall(r'`topic:([^`]+)`', tags))
            risk_type = "E" if "`type:E`" in tags else ("A" if "`type:A`" in tags else "")
            papers.append(dict(
                title=title.strip(), authors=authors.strip(), year="20"+arxiv[:2],
                venue=venue.strip(), arxiv=arxiv.strip(), section=sec_title,
                rf=",".join(sorted(set(rfs))).upper() or rf,
                risk_type=risk_type, topics=topics,
                notes=notes.strip().strip('—').strip(), tags=tags.strip(),
            ))
    return papers

PAPERS = parse_readme()

RF_ACCENT = {
    'RF0': RGBColor(0x2D,0x6A,0x9F), 'RF1': RGBColor(0x2D,0x6A,0x9F),
    'RF2': RGBColor(0xE0,0x7B,0x39), 'RF3': RGBColor(0x3A,0x9A,0x5B),
    'RF4': RGBColor(0x8E,0x44,0xAD), 'RF5': RGBColor(0xC0,0x39,0x2B),
    'RF6': RGBColor(0x17,0x7E,0x89), 'EVAL': RGBColor(0x55,0x55,0x55),
    'DEF': RGBColor(0x3A,0x9A,0x5B), 'SURVEY': RGBColor(0x1F,0x39,0x64),
    'TBD': RGBColor(0x88,0x88,0x88), 'OTHER': RGBColor(0x88,0x88,0x88),
}
DARK  = RGBColor(0x22,0x22,0x22)
MID   = RGBColor(0x44,0x44,0x44)
LIGHT = RGBColor(0x88,0x88,0x88)
WHITE = RGBColor(0xFF,0xFF,0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

def txt(slide, l, t, w, h, text, size=10, bold=False, color=DARK,
        align_h="left", italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT if align_h == "left" else PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color

def bullets(slide, l, t, w, h, lines, size=9, color=DARK):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = "• " + line
        run.font.size = Pt(size); run.font.color.rgb = color

def hline(slide, l, t, w, rgb, thickness=Inches(0.03)):
    s = slide.shapes.add_shape(1, l, t, w, thickness)
    s.fill.solid(); s.fill.fore_color.rgb = rgb
    s.line.fill.background()

current_section = None
for p in PAPERS:
    rf = p['rf'].split(',')[0]
    accent = RF_ACCENT.get(rf, RGBColor(0x55,0x55,0x55))
    k = KNOWN.get(p['arxiv'], {})

    # section divider
    if p['section'] != current_section:
        current_section = p['section']
        ds = prs.slides.add_slide(prs.slide_layouts[6])
        hline(ds, Inches(0.5), Inches(3.3), Inches(12.3), accent, Inches(0.06))
        txt(ds, Inches(0.5), Inches(3.45), Inches(12.3), Inches(1.0),
            p['section'], size=34, bold=True, color=accent, align_h="center")
        txt(ds, Inches(0.5), Inches(4.55), Inches(12.3), Inches(0.5),
            f"Risk Family: {rf}", size=18, color=LIGHT, align_h="center")

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # RF tag (no accent bar — just text)
    rf_label = rf + (f"  ·  Type {p['risk_type']}" if p['risk_type'] else "")
    txt(slide, Inches(0.3), Inches(0.18), Inches(4), Inches(0.32),
        rf_label, size=11, bold=True, color=accent)

    # title
    txt(slide, Inches(0.3), Inches(0.48), Inches(12.7), Inches(0.8),
        p['title'], size=17, bold=True, color=DARK)

    # authors / venue / arxiv
    txt(slide, Inches(0.3), Inches(1.25), Inches(12), Inches(0.3),
        f"{p['authors']}  ·  {p['venue']}  ·  {p['year']}  ·  arxiv.org/abs/{p['arxiv']}",
        size=11, color=LIGHT, italic=True)

    hline(slide, Inches(0.3), Inches(1.58), Inches(12.7), RGBColor(0xDD,0xDD,0xDD))

    # ── left column ───────────────────────────────────────────────────────────
    cl, cw = Inches(0.3), Inches(5.9)
    y = Inches(1.72)

    domain = k.get('domain', '')
    if domain:
        txt(slide, cl, y, cw, Inches(0.28), "Domain", size=10, bold=True, color=LIGHT)
        txt(slide, cl, y+Inches(0.26), cw, Inches(0.32), domain, size=12, color=DARK)
        y += Inches(0.65)

    metrics = k.get('metrics', [])
    txt(slide, cl, y, cw, Inches(0.28), "Metrics", size=10, bold=True, color=LIGHT)
    y += Inches(0.28)
    if metrics:
        mh = Inches(0.26 * len(metrics) + 0.1)
        bullets(slide, cl, y, cw, mh, metrics, size=11)
        y += mh + Inches(0.2)
    else:
        txt(slide, cl, y, cw, Inches(0.3), "(not yet extracted)", size=11, color=LIGHT, italic=True)
        y += Inches(0.45)

    datasets = k.get('datasets', [])
    txt(slide, cl, y, cw, Inches(0.28), "Datasets / Benchmarks", size=10, bold=True, color=LIGHT)
    y += Inches(0.28)
    if datasets:
        bullets(slide, cl, y, cw, Inches(0.26*len(datasets)+0.1), datasets, size=11)
    else:
        txt(slide, cl, y, cw, Inches(0.3), "(not yet extracted)", size=11, color=LIGHT, italic=True)

    # ── right column ──────────────────────────────────────────────────────────
    cr, crw = Inches(6.4), Inches(6.6)
    yr = Inches(1.72)

    highlights = k.get('highlights', [])
    txt(slide, cr, yr, crw, Inches(0.28), "Contribution & Highlights", size=10, bold=True, color=LIGHT)
    yr += Inches(0.28)
    if highlights:
        bullets(slide, cr, yr, crw, H-yr-Inches(0.2), highlights, size=12)
    else:
        fallback = []
        if p['topics']:
            fallback.append(f"Topics: {p['topics']}")
        if p['notes']:
            fallback.append(f"Notes: {p['notes'][:200]}")
        fallback.append(f"Tags: {p['tags'][:150]}")
        bullets(slide, cr, yr, crw, H-yr-Inches(0.2), fallback, size=11, color=MID)

prs.save("paper_slides.pptx")
print(f"Saved → paper_slides.pptx  ({len(PAPERS)} papers)")
