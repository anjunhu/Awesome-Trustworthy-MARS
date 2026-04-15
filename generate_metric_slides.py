"""
Generate evaluation-metric slides for the MA-RS reading list.

Slide 1 – By Provenance  : one row per paper, columns = paper, domain, metrics, datasets
Slide 2 – By Risk Family : one section per RF, listing metrics + sources

Run:  python3 generate_metric_slides.py
Output: metric_slides.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import textwrap

# ── colour palette ────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1F, 0x39, 0x64)
MID_BLUE   = RGBColor(0x2D, 0x6A, 0x9F)
ORANGE     = RGBColor(0xE0, 0x7B, 0x39)
GREEN      = RGBColor(0x3A, 0x9A, 0x5B)
LIGHT_GREY = RGBColor(0xF4, 0xF6, 0xF8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)

# ── paper data ────────────────────────────────────────────────────────────────
# Each entry: (short_name, risk_families, domain, metrics, datasets)
PAPERS = [
    # RF1 – Prompt Injection & Jailbreaking
    (
        "InjecAgent\n(Zhan et al., ACL'24)",
        ["RF1"],
        "General agent\n(tool-use)",
        "ASR-valid: ASR among valid outputs\n"
        "ASR-all: ASR across all outputs\n"
        "Sensitivity Rate: % outputs recognising attack as abnormal\n"
        "Valid Rate: % outputs parseable as valid agent actions",
        "InjecAgent benchmark\n(1,054 test cases, 17 user tools, 62 attacker tools)",
    ),
    (
        "LlamaFirewall\n(Meta AI, 2025)",
        ["RF1"],
        "General agent\n(coding/tool-use)",
        "ASR: Attack Success Rate (post-guardrail)\n"
        "Utility: Task success rate under defence\n"
        "AUC: Area Under ROC for jailbreak detection\n"
        "Recall @ 1% FPR: detection rate at low false-positive threshold\n"
        "Precision / Recall: CodeShield insecure-code detection",
        "AgentDojo (97 tasks)\n"
        "Meta internal goal-hijacking benchmark (600 scenarios)\n"
        "CyberSecEval3 (50 code completions/language)",
    ),
    (
        "CORBA\n(Zhou et al., 2025)",
        ["RF1", "RF6"],
        "General MAS\n(AutoGen, Camel)",
        "P-ASR: Proportional Attack Success Rate\n"
        "  = fraction of agents entering blocked state\n"
        "PTN: Peak Blocking Turn Number\n"
        "  = turns until attack stabilises at max P-ASR\n"
        "PPL: Perplexity of attack prompt (stealthiness proxy)",
        "AutoGen & Camel frameworks\n"
        "Open-ended MAS (6-agent free dialogue)\n"
        "Models: GPT-4o-mini, GPT-4, GPT-3.5-turbo, Gemini-2.0-Flash,\n"
        "  Qwen2.5-14B, Llama3.1-70B, Gemma-2-27B",
    ),
    (
        "AiTM\n(He et al., ACL'25)",
        ["RF3"],
        "General MAS\n(coding/QA)",
        "ASR: Attack Success Rate\n"
        "  (targeted behaviour or DoS success fraction)",
        "MMLU (biology, physics)\n"
        "HumanEval (164 coding problems)\n"
        "MBPP (974 coding tasks)\n"
        "SoftwareDev (MetaGPT tasks)\n"
        "Frameworks: AutoGen, Camel, MetaGPT, ChatDev",
    ),
    # RF2 – Data Poisoning & Backdoor
    (
        "DrunkAgent\n(Yang et al., 2025)",
        ["RF2"],
        "RecSys\n(agentic CF/RAG/Seq)",
        "HR@K: Hit Ratio at K (attack transferability)\n"
        "NDCG@K: Normalised Discounted Cumulative Gain at K\n"
        "Perplexity ↓: text naturalness (stealthiness)\n"
        "Strategy Success Rate: % target agents 'drunk'",
        "Amazon Review Data:\n"
        "  CDs & Vinyl, Office Products, Musical Instruments\n"
        "Yelp (cross-domain generalisation)\n"
        "Victim systems: AgentCF, AgentRAG, AgentSEQ",
    ),
    (
        "PeerGuard\n(Fan & Li, 2025)",
        ["RF2"],
        "General MAS\n(QA/reasoning)",
        "TPR: True Positive Rate (backdoor-triggered inputs detected)\n"
        "FPR: False Positive Rate (clean inputs incorrectly flagged)",
        "MMLU (57 fields)\n"
        "CSQA (12,247 common-sense questions)\n"
        "ARC-Easy / ARC-Challenge\n"
        "Models: GPT-4o, Llama3-70B\n"
        "Frameworks: AutoGen, Camel",
    ),
    # RF3 – Inter-Agent Communication
    (
        "Compositional Privacy\n(Patil et al., 2025)",
        ["RF3", "RF4"],
        "General MAS\n(enterprise/org)",
        "Leakage Accuracy: I[ŝ = s*] — adversary infers sensitive target\n"
        "Sensitive Blocked (%): sensitive scenarios fully blocked\n"
        "Benign Succeeded (%): benign queries answered correctly\n"
        "Balanced Outcome: avg(Sensitive Blocked, Benign Succeeded)\n"
        "Overall Success: both blocked & benign succeed simultaneously\n"
        "PlanExec@m: all m plan steps executed correctly",
        "Synthetic multi-agent scenarios (119 paired adv/benign)\n"
        "Models: Qwen3-32B, Gemini-2.5-Pro, GPT-5",
    ),
    # RF4 – Privacy & Inversion
    (
        "TrojanStego\n(Anon., 2025)",
        ["RF4"],
        "General LLM\n(steganography)",
        "Steganographic capacity (bits/token)\n"
        "Detection evasion rate vs. statistical detectors\n"
        "Task utility preservation (downstream accuracy)",
        "Standard NLP benchmarks\n"
        "(paper-specific; details in full text)",
    ),
    # RF5 – Cognitive Bias & Dark Patterns
    (
        "DarkBench\n(Kran et al., 2025)",
        ["RF5"],
        "General chatbot\n(LLM products)",
        "Dark Pattern Rate (%): fraction of conversations exhibiting\n"
        "  each of 6 categories:\n"
        "  Brand Bias, User Retention, Sycophancy,\n"
        "  Anthropomorphism, Harmful Generation, Sneaking\n"
        "Cohen's Kappa / Jaccard / Agreement Rate: inter-annotator agreement",
        "DarkBench (660 adversarial prompts × 14 models)\n"
        "Models: GPT-3.5/4/4-Turbo/4o, Claude-3 family,\n"
        "  Gemini-1.0/1.5, Llama-3-8B/70B, Mistral-7B, Mixtral-8x7B",
    ),
    (
        "Bias Beware\n(Krasniqi et al., EMNLP'25)",
        ["RF5"],
        "RecSys\n(product rec.)",
        "Cognitive Bias Score: LLM-judged presence of\n"
        "  anchoring, availability, confirmation bias\n"
        "Recommendation Accuracy: standard RecSys metrics\n"
        "Bias-Accuracy Trade-off",
        "Amazon product recommendation scenarios\n"
        "(paper-specific datasets)",
    ),
    # RF6 – Availability / Collusion
    (
        "Emergent Social Intelligence\n(Huang et al., 2026)",
        ["RF6"],
        "General MAS\n(social simulation)",
        "Collusion Rate: fraction of agent pairs exhibiting\n"
        "  coordinated unsafe behaviour\n"
        "Social Intelligence Score: emergent cooperation metrics\n"
        "Safety Violation Rate: policy-breaking outputs",
        "RiskLab benchmark\n"
        "(decentralised MAS social simulation)",
    ),
]

# ── risk family metadata ──────────────────────────────────────────────────────
RISK_FAMILIES = [
    {
        "id": "RF1",
        "name": "Prompt Injection & Jailbreaking",
        "type": "Amplified + Emergent",
        "color": MID_BLUE,
        "metrics": [
            ("ASR / ASR-valid / ASR-all",
             "% trials achieving attacker goal; -valid restricts to parseable outputs",
             "InjecAgent, AiTM, LlamaFirewall, CORBA"),
            ("Sensitivity Rate",
             "% agent outputs recognising injected instruction as abnormal",
             "InjecAgent"),
            ("Valid Rate",
             "% outputs parseable as valid agent actions (proxy for model capability)",
             "InjecAgent"),
            ("AUC / Recall@1%FPR",
             "Classifier performance for jailbreak detection at low false-positive rate",
             "LlamaFirewall (PromptGuard 2)"),
            ("Utility",
             "Task success rate under defence (measures defence cost)",
             "LlamaFirewall, AiTM"),
        ],
    },
    {
        "id": "RF2",
        "name": "Data Poisoning & Backdoor Attacks",
        "type": "Amplified",
        "color": ORANGE,
        "metrics": [
            ("HR@K / NDCG@K",
             "Hit Ratio / Normalised DCG at rank K — attack transferability in RecSys",
             "DrunkAgent"),
            ("Perplexity ↓",
             "Text naturalness of adversarial item description (lower = stealthier)",
             "DrunkAgent"),
            ("Strategy Success Rate",
             "% target item agents successfully 'drunk' (memory corrupted)",
             "DrunkAgent"),
            ("TPR / FPR",
             "True/False Positive Rate for backdoor detection",
             "PeerGuard"),
            ("ASR (backdoor)",
             "% triggered inputs producing target (wrong) output",
             "PeerGuard"),
        ],
    },
    {
        "id": "RF3",
        "name": "Inter-Agent Communication Attacks",
        "type": "Emergent",
        "color": GREEN,
        "metrics": [
            ("ASR (communication)",
             "% inter-agent message manipulations achieving attacker goal",
             "AiTM, CORBA"),
            ("P-ASR",
             "Proportional ASR = fraction of agents entering blocked/compromised state",
             "CORBA"),
            ("PTN",
             "Peak Blocking Turn Number = turns until attack stabilises at max P-ASR",
             "CORBA"),
            ("Leakage Accuracy",
             "I[ŝ = s*]: adversary correctly infers sensitive target via composition",
             "Compositional Privacy"),
            ("PlanExec@m",
             "All m adversarial plan steps executed correctly (execution fidelity)",
             "Compositional Privacy"),
        ],
    },
    {
        "id": "RF4",
        "name": "Privacy & Inversion Attacks",
        "type": "Amplified + Emergent",
        "color": RGBColor(0x8E, 0x44, 0xAD),
        "metrics": [
            ("Sensitive Blocked (%)",
             "% sensitive scenarios where adversary completely fails to infer target",
             "Compositional Privacy"),
            ("Benign Succeeded (%)",
             "% benign queries answered correctly without obstruction",
             "Compositional Privacy"),
            ("Balanced Outcome",
             "avg(Sensitive Blocked, Benign Succeeded) — privacy-utility trade-off",
             "Compositional Privacy"),
            ("Steganographic Capacity",
             "Bits of covert information per token (TrojanStego)",
             "TrojanStego"),
            ("Leakage Rate / τ_leak",
             "Fraction of PII entities recovered; time to first leak",
             "MAMA (from provenance table)"),
        ],
    },
    {
        "id": "RF5",
        "name": "Cognitive Bias & Dark Patterns",
        "type": "Amplified",
        "color": RGBColor(0xC0, 0x39, 0x2B),
        "metrics": [
            ("Dark Pattern Rate (%)",
             "Fraction of conversations exhibiting each of 6 dark pattern categories",
             "DarkBench"),
            ("Cohen's Kappa / Jaccard",
             "Inter-annotator agreement for dark pattern annotation",
             "DarkBench"),
            ("Cognitive Bias Score",
             "LLM-judged presence of anchoring, availability, confirmation bias",
             "Bias Beware"),
            ("ASR-H@r / ASR-N@r",
             "Drop in Hit Ratio / NDCG caused by bias-inducing attack",
             "CheatAgent (from provenance table)"),
        ],
    },
    {
        "id": "RF6",
        "name": "Availability, Collusion & Resource Depletion",
        "type": "Emergent",
        "color": RGBColor(0x17, 0x7E, 0x89),
        "metrics": [
            ("P-ASR / PTN",
             "Proportional ASR and Peak Blocking Turn Number (availability attacks)",
             "CORBA"),
            ("Collusion Rate",
             "Fraction of agent pairs exhibiting coordinated unsafe behaviour",
             "Emergent Social Intelligence"),
            ("Safety Violation Rate",
             "% outputs violating safety policy under collusion",
             "Emergent Social Intelligence"),
            ("Agent Blocking Rate",
             "Fraction of agents rendered non-functional",
             "CORBA"),
            ("CPU/Memory/Latency Delta",
             "Computational overhead induced by attack/defence",
             "TOMA (from provenance table)"),
        ],
    },
]

# ── helpers ───────────────────────────────────────────────────────────────────

def set_cell_bg(cell, rgb: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")


def add_cell_text(cell, text, font_size=9, bold=False, color=BLACK, wrap=True):
    tf = cell.text_frame
    tf.word_wrap = wrap
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    w, h = prs.slide_width, prs.slide_height

    # background
    bg = slide.shapes.add_shape(1, 0, 0, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # title
    txb = slide.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(12), Inches(1.5))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Evaluation Metrics in MA-RS Security Research"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # subtitle
    txb2 = slide.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(12), Inches(1))
    tf2 = txb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "By Provenance  ·  By Risk Family (Kurt's Taxonomy)"
    run2.font.size = Pt(18)
    run2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xFF)


def add_provenance_slide(prs):
    """One wide table: Paper | Domain | Metrics | Datasets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w, h = prs.slide_width, prs.slide_height

    # header bar
    hdr = slide.shapes.add_shape(1, 0, 0, w, Inches(0.65))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = DARK_BLUE
    hdr.line.fill.background()

    txb = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(12), Inches(0.5))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Evaluation Framework  ·  What to Measure (By Provenance)"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # table geometry
    rows = len(PAPERS) + 1  # +1 header
    cols = 4
    left, top = Inches(0.15), Inches(0.72)
    tbl_w = w - Inches(0.3)
    tbl_h = h - Inches(0.85)

    tbl = slide.shapes.add_table(rows, cols, left, top, tbl_w, tbl_h).table

    # column widths (proportional)
    col_fracs = [0.18, 0.10, 0.42, 0.30]
    for i, frac in enumerate(col_fracs):
        tbl.columns[i].width = int(tbl_w * frac)

    # header row
    headers = ["Paper", "Domain", "Metrics (Acronym: Definition)", "Datasets"]
    for ci, hdr_txt in enumerate(headers):
        cell = tbl.cell(0, ci)
        set_cell_bg(cell, DARK_BLUE)
        add_cell_text(cell, hdr_txt, font_size=10, bold=True, color=WHITE)

    # data rows
    for ri, (name, rfs, domain, metrics, datasets) in enumerate(PAPERS, start=1):
        row_bg = LIGHT_GREY if ri % 2 == 0 else WHITE
        rf_tag = ", ".join(rfs)

        cell_name = tbl.cell(ri, 0)
        set_cell_bg(cell_name, row_bg)
        add_cell_text(cell_name, f"{name}\n[{rf_tag}]", font_size=8, bold=False)

        cell_dom = tbl.cell(ri, 1)
        set_cell_bg(cell_dom, row_bg)
        add_cell_text(cell_dom, domain, font_size=8)

        cell_met = tbl.cell(ri, 2)
        set_cell_bg(cell_met, row_bg)
        add_cell_text(cell_met, metrics, font_size=7.5)

        cell_dat = tbl.cell(ri, 3)
        set_cell_bg(cell_dat, row_bg)
        add_cell_text(cell_dat, datasets, font_size=7.5)


def add_risk_family_slide(prs):
    """One wide table: RF | Metric | Definition | Source papers"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w, h = prs.slide_width, prs.slide_height

    # header bar
    hdr = slide.shapes.add_shape(1, 0, 0, w, Inches(0.65))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = DARK_BLUE
    hdr.line.fill.background()

    txb = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(12), Inches(0.5))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Evaluation Framework  ·  What to Measure (By Risk Family)"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # flatten rows: (rf_label, metric, definition, sources)
    flat_rows = []
    for rf in RISK_FAMILIES:
        for i, (metric, defn, sources) in enumerate(rf["metrics"]):
            rf_label = f"{rf['id']}\n{rf['name']}\n({rf['type']})" if i == 0 else ""
            flat_rows.append((rf_label, rf["color"], metric, defn, sources))

    rows = len(flat_rows) + 1
    cols = 4
    left, top = Inches(0.15), Inches(0.72)
    tbl_w = w - Inches(0.3)
    tbl_h = h - Inches(0.85)

    tbl = slide.shapes.add_table(rows, cols, left, top, tbl_w, tbl_h).table

    col_fracs = [0.18, 0.22, 0.38, 0.22]
    for i, frac in enumerate(col_fracs):
        tbl.columns[i].width = int(tbl_w * frac)

    # header
    for ci, hdr_txt in enumerate(
        ["Risk Family", "Metric", "Definition / Formula", "Source Papers"]
    ):
        cell = tbl.cell(0, ci)
        set_cell_bg(cell, DARK_BLUE)
        add_cell_text(cell, hdr_txt, font_size=10, bold=True, color=WHITE)

    # data
    for ri, (rf_label, rf_color, metric, defn, sources) in enumerate(flat_rows, start=1):
        cell_rf = tbl.cell(ri, 0)
        set_cell_bg(cell_rf, rf_color)
        add_cell_text(cell_rf, rf_label, font_size=8, bold=bool(rf_label), color=WHITE)

        row_bg = LIGHT_GREY if ri % 2 == 0 else WHITE

        cell_met = tbl.cell(ri, 1)
        set_cell_bg(cell_met, row_bg)
        add_cell_text(cell_met, metric, font_size=8, bold=True)

        cell_def = tbl.cell(ri, 2)
        set_cell_bg(cell_def, row_bg)
        add_cell_text(cell_def, defn, font_size=7.5)

        cell_src = tbl.cell(ri, 3)
        set_cell_bg(cell_src, row_bg)
        add_cell_text(cell_src, sources, font_size=7.5)


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)   # widescreen 16:9
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_provenance_slide(prs)
    add_risk_family_slide(prs)

    out = "metric_slides.pptx"
    prs.save(out)
    print(f"Saved → {out}")


if __name__ == "__main__":
    main()
